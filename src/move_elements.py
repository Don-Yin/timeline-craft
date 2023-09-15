import collections
import collections.abc
from pptx.enum.shapes import MSO_SHAPE


from src.utils.shapes import set_shape_transparency, send_backwards
from src.utils.placeholder import add_placeholder
from src.utils.paragraph import add_paragraph, amend_font


def move_elements_to_right(ppt, sidebar_width=0.12):
    """
    Moving all placeholders and images to the right of the screen (the content area)
    Note: later reduce the font size
    Note: the bullet pint issue is not caused by this function.
    """
    for slide in ppt.slides:
        for shape in slide.shapes:
            top, height, left, width = shape.top, shape.height, shape.left, shape.width

            right = left + width
            # bottom = top + height

            loc_left = left / ppt.slide_width
            loc_right = right / ppt.slide_width
            # loc_top = top / ppt.slide_width
            # loc_bottom = bottom / ppt.slide_width

            content_space_width = ppt.slide_width * (1 - sidebar_width)
            new_left = ppt.slide_width * sidebar_width + (
                loc_left * content_space_width
            )
            new_right = ppt.slide_width * sidebar_width + (
                loc_right * content_space_width
            )
            new_width = new_right - new_left

            shape.left, shape.width = int(new_left), int(new_width)
            shape.top, shape.height = int(top), int(height)


def merge_tags(tags: list[str]) -> list[str]:
    """
    Merge the adjacent slides with the same tag
    """
    merged_tags = []
    for i, tag in enumerate(tags):
        if i == 0:
            merged_tags.append(tag)
        elif tag != merged_tags[-1]:
            merged_tags.append(tag)
    return merged_tags


def set_sidebar_timeline(
    ppt,
    tags: list[int],
    sidebar_width,
    sidebar_transparency,
    sidebar_color,
    sidebar_color_outline,
    sidebar_item_height,
    sidebar_init_font_size,
    sidebar_item_font,
    sidebar_item_font_color,
    indicator_color,
    indicator_transparency,
):
    assert len(tags) == len(
        ppt.slides
    ), f"The number of tags: {len(tags)} has to match the number of slides: {len(ppt.slides)}"

    # merge the adjacent slides with the same tag
    merged_tags = merge_tags(tags)

    # adding the base shapes to each slide
    for slide in ppt.slides:
        # ------------------------ shaping the sidebar itself ------------------------
        sidebar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, ppt.slide_width * sidebar_width, ppt.slide_height
        )
        sidebar.name = "!!SIDEBAR"
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = sidebar_color
        sidebar.line.color.rgb = sidebar_color_outline
        set_shape_transparency(sidebar, sidebar_transparency)

        offset = 0  # sidebar item offset from top
        for tag in merged_tags:
            placeholder = add_placeholder(
                ppt=ppt,
                slide_index=ppt.slides.index(slide),
                template="FOOTER",
                left=0,
                top=offset,
                width=sidebar_width,
                height=sidebar_item_height,
            )

            add_paragraph(
                placeholder=placeholder,
                text=tag,
                font_size=sidebar_init_font_size,
                font_family=sidebar_item_font,
                font_color=sidebar_item_font_color,
            )

            setattr(placeholder, "name", f"!!SIDEBAR_{merged_tags.index(tag)}")

            if tags[ppt.slides.index(slide)] == tag:
                amend_font(
                    placeholder=placeholder,
                    font_family=sidebar_item_font,
                    font_size=sidebar_init_font_size + 3,
                    bold=True,
                )

                indicator = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left=0,
                    top=offset * ppt.slide_height,
                    width=ppt.slide_width * (sidebar_width + 0.01),
                    height=ppt.slide_height * sidebar_item_height,
                )

                indicator.name = "!!INDICATOR"
                indicator.fill.solid()
                indicator.fill.fore_color.rgb = indicator_color
                indicator.line.color.rgb = sidebar_color_outline
                set_shape_transparency(sidebar, indicator_transparency)
                send_backwards(slide, indicator)

            offset += sidebar_item_height

        send_backwards(slide, sidebar)
