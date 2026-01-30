"""
Preview service with persistent LibreOffice (unoserver) for fast PPTX rendering.
Provides SSE progress updates during rendering.
"""

import asyncio
import base64
import io
import json
import logging
import os
import tempfile
import time
from pathlib import Path
from typing import AsyncGenerator

import fitz  # PyMuPDF
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from minio import Minio
from pydantic import BaseModel, Field
from subprocess import run, PIPE, CalledProcessError

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="Preview Service",
    description="Fast PPTX rendering with persistent LibreOffice (unoserver)",
    version="1.0.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# MinIO configuration
MINIO_ENDPOINT = os.getenv("MINIO_ENDPOINT", "file-store:9000")
MINIO_ACCESS_KEY = os.getenv("MINIO_ACCESS_KEY", "minioadmin")
MINIO_SECRET_KEY = os.getenv("MINIO_SECRET_KEY", "minioadmin")
MINIO_BUCKET_UPLOADS = os.getenv("MINIO_BUCKET_UPLOADS", "uploads")

minio_client = Minio(MINIO_ENDPOINT, access_key=MINIO_ACCESS_KEY, secret_key=MINIO_SECRET_KEY, secure=False)


class PreviewRequest(BaseModel):
    """Request for rendering preview thumbnails"""

    tags: list[str] = Field(description="List of section tags, one per slide")
    sidebar_width: float = Field(default=0.12, ge=0.05, le=0.5)
    sidebar_item_height: float = Field(default=0.10, ge=0.03, le=0.3)


def get_file_from_minio(file_id: str) -> bytes:
    """Retrieve file from MinIO"""
    try:
        response = minio_client.get_object(MINIO_BUCKET_UPLOADS, file_id)
        data = response.read()
        response.close()
        response.release_conn()
        return data
    except Exception as e:
        raise HTTPException(status_code=404, detail=f"File not found: {str(e)}")


def convert_pptx_to_png_first_slide(pptx_bytes: bytes, job_id: str) -> bytes:
    """convert pptx directly to png (first slide only) - 27x faster than pdf!"""
    with tempfile.TemporaryDirectory(prefix=f"preview-{job_id}-") as tmpdir:
        tmpdir_path = Path(tmpdir)
        pptx_path = tmpdir_path / "source.pptx"
        pptx_path.write_bytes(pptx_bytes)

        cmd = [
            "soffice",
            "--headless",
            f"-env:UserInstallation=file://{str(tmpdir_path)}/profile",
            "--convert-to",
            "png",
            "--outdir",
            str(tmpdir_path),
            str(pptx_path),
        ]

        run(cmd, check=True, stdout=PIPE, stderr=PIPE, timeout=60)
        logger.info(f"direct png export completed for {job_id}")

        png_path = tmpdir_path / "source.png"
        if not png_path.exists():
            raise HTTPException(status_code=500, detail="PNG export failed")

        return png_path.read_bytes()


def convert_pptx_to_pdf(pptx_bytes: bytes, job_id: str) -> bytes:
    """convert pptx to pdf using libreoffice (slower, all slides)"""
    with tempfile.TemporaryDirectory(prefix=f"preview-{job_id}-") as tmpdir:
        tmpdir_path = Path(tmpdir)
        pptx_path = tmpdir_path / "source.pptx"

        pptx_path.write_bytes(pptx_bytes)

        # Use soffice with isolated user profile for parallel execution
        cmd = [
            "soffice",
            "--headless",
            f"-env:UserInstallation=file://{str(tmpdir_path)}/profile",
            "--convert-to",
            "pdf",
            "--outdir",
            str(tmpdir_path),
            str(pptx_path),
        ]

        try:
            result = run(cmd, check=True, stdout=PIPE, stderr=PIPE, timeout=180)
            logger.info(f"soffice conversion completed for {job_id}")
        except CalledProcessError as exc:
            logger.error(f"soffice failed: {exc.stderr.decode(errors='ignore')}")
            raise HTTPException(status_code=500, detail="Failed to convert PPTX to PDF")
        except Exception as exc:
            logger.error(f"soffice error: {exc}")
            raise HTTPException(status_code=500, detail=f"Conversion error: {str(exc)}")

        pdf_path = pptx_path.with_suffix(".pdf")
        if not pdf_path.exists():
            raise HTTPException(status_code=500, detail="Converted PDF not found")

        return pdf_path.read_bytes()


def render_pdf_page_to_image(pdf_bytes: bytes, page_index: int, target_width: int = 640) -> str:
    """Render a single PDF page to base64 PNG"""
    doc = fitz.open("pdf", pdf_bytes)
    if page_index < 0 or page_index >= doc.page_count:
        doc.close()
        raise HTTPException(status_code=404, detail="Page index out of range")

    page = doc.load_page(page_index)
    scale = target_width / (page.rect.width or 1)
    pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
    img_bytes = pix.tobytes("jpeg", jpg_quality=80)
    doc.close()

    return base64.b64encode(img_bytes).decode("utf-8")


async def generate_preview_sse(file_id: str, pptx_bytes: bytes) -> AsyncGenerator[str, None]:
    """
    Generate preview images with SSE progress updates.
    Yields JSON events with progress info.
    """
    job_id = file_id[:8]
    total_stages = 3  # convert, render, done

    try:
        # Stage 1: Converting PPTX to PDF
        yield f"data: {json.dumps({'stage': 'converting', 'progress': 0, 'message': 'converting pptx to pdf...'})}\n\n"

        start_convert = time.time()
        pdf_bytes = convert_pptx_to_pdf(pptx_bytes, job_id)
        convert_time = time.time() - start_convert

        yield f"data: {json.dumps({'stage': 'converting', 'progress': 33, 'message': f'converted in {convert_time:.1f}s'})}\n\n"

        # Stage 2: Rendering PDF pages to images
        doc = fitz.open("pdf", pdf_bytes)
        total_pages = doc.page_count
        doc.close()

        yield f"data: {json.dumps({'stage': 'rendering', 'progress': 33, 'message': f'rendering {total_pages} slides...', 'total_slides': total_pages})}\n\n"

        thumbnails = []
        start_render = time.time()

        for i in range(total_pages):
            thumbnail = render_pdf_page_to_image(pdf_bytes, i)
            thumbnails.append(thumbnail)

            # Calculate progress (33-95%)
            render_progress = 33 + int((i + 1) / total_pages * 62)

            yield f"data: {json.dumps({'stage': 'rendering', 'progress': render_progress, 'message': f'rendered slide {i+1}/{total_pages}', 'current_slide': i+1, 'total_slides': total_pages})}\n\n"

            # Small delay to allow client to process
            await asyncio.sleep(0.01)

        render_time = time.time() - start_render

        # Stage 3: Done
        yield f"data: {json.dumps({'stage': 'done', 'progress': 100, 'message': f'completed in {convert_time + render_time:.1f}s', 'thumbnails': thumbnails, 'format': 'jpeg'})}\n\n"

    except HTTPException as e:
        yield f"data: {json.dumps({'stage': 'error', 'progress': 0, 'message': e.detail})}\n\n"
    except Exception as e:
        logger.error(f"Preview generation error: {e}")
        yield f"data: {json.dumps({'stage': 'error', 'progress': 0, 'message': str(e)})}\n\n"


@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {"status": "healthy", "service": "preview"}


@app.get("/render-previews/{file_id}")
async def render_previews_sse(file_id: str):
    """
    Render all slide previews with SSE progress updates.

    Returns a stream of JSON events:
    - {stage: 'converting', progress: 0-33, message: '...'}
    - {stage: 'rendering', progress: 33-95, message: '...', current_slide: N, total_slides: M}
    - {stage: 'done', progress: 100, thumbnails: [...], format: 'jpeg'}
    - {stage: 'error', message: '...'}
    """
    logger.info(f"Starting preview render for {file_id}")

    # Get file from MinIO
    pptx_bytes = get_file_from_minio(file_id)

    return StreamingResponse(
        generate_preview_sse(file_id, pptx_bytes),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/render-first-slide/{file_id}")
async def render_first_slide_sse(file_id: str, request: PreviewRequest):
    """render only the first slide with sidebar - uses fast direct png export (~1-3s)"""
    logger.info(f"rendering first slide preview for {file_id}")

    pptx_bytes = get_file_from_minio(file_id)

    async def generate():
        job_id = file_id[:8]
        start_time = time.time()

        yield f"data: {json.dumps({'stage': 'processing', 'progress': 0, 'message': 'applying sidebar...'})}\n\n"

        from timeline import Configurations, move_elements_to_right, set_sidebar_timeline
        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))

        if len(request.tags) != len(prs.slides):
            yield f"data: {json.dumps({'stage': 'error', 'message': f'tags count ({len(request.tags)}) must match slides ({len(prs.slides)})'})}\n\n"
            return

        yield f"data: {json.dumps({'stage': 'processing', 'progress': 30, 'message': 'processing presentation...'})}\n\n"

        config = Configurations(sidebar_width=request.sidebar_width, sidebar_item_height=request.sidebar_item_height)
        move_elements_to_right(prs, config=config)
        set_sidebar_timeline(ppt=prs, tags=request.tags, config=config)

        processed_buffer = io.BytesIO()
        prs.save(processed_buffer)
        processed_bytes = processed_buffer.getvalue()

        yield f"data: {json.dumps({'stage': 'rendering', 'progress': 60, 'message': 'rendering first slide...'})}\n\n"

        # Use fast direct PNG export (only renders first slide, ~1s instead of ~28s)
        png_bytes = convert_pptx_to_png_first_slide(processed_bytes, job_id)
        thumbnail = base64.b64encode(png_bytes).decode("utf-8")

        total_time = time.time() - start_time
        logger.info(f"first slide preview completed in {total_time:.1f}s")

        yield f"data: {json.dumps({'stage': 'done', 'progress': 100, 'message': f'completed in {total_time:.1f}s', 'thumbnails': [thumbnail], 'format': 'png'})}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive", "X-Accel-Buffering": "no"},
    )


@app.post("/render-previews-with-sidebar/{file_id}")
async def render_previews_with_sidebar_sse(file_id: str, request: PreviewRequest):
    """render all slides with sidebar applied"""
    logger.info(f"starting preview render with sidebar for {file_id}")

    pptx_bytes = get_file_from_minio(file_id)

    async def generate():
        job_id = file_id[:8]

        yield f"data: {json.dumps({'stage': 'processing', 'progress': 0, 'message': 'applying sidebar to slides...'})}\n\n"

        from timeline import Configurations, move_elements_to_right, set_sidebar_timeline
        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))

        if len(request.tags) != len(prs.slides):
            yield f"data: {json.dumps({'stage': 'error', 'message': f'tags count ({len(request.tags)}) must match slides ({len(prs.slides)})'})}\n\n"
            return

        yield f"data: {json.dumps({'stage': 'processing', 'progress': 10, 'message': 'moving elements...'})}\n\n"

        config = Configurations(sidebar_width=request.sidebar_width, sidebar_item_height=request.sidebar_item_height)
        move_elements_to_right(prs, config=config)

        yield f"data: {json.dumps({'stage': 'processing', 'progress': 20, 'message': 'adding sidebar timeline...'})}\n\n"

        set_sidebar_timeline(ppt=prs, tags=request.tags, config=config)

        processed_buffer = io.BytesIO()
        prs.save(processed_buffer)
        processed_bytes = processed_buffer.getvalue()

        yield f"data: {json.dumps({'stage': 'converting', 'progress': 25, 'message': 'converting to pdf...'})}\n\n"

        start_convert = time.time()
        pdf_bytes = convert_pptx_to_pdf(processed_bytes, job_id)
        convert_time = time.time() - start_convert

        yield f"data: {json.dumps({'stage': 'converting', 'progress': 50, 'message': f'converted in {convert_time:.1f}s'})}\n\n"

        doc = fitz.open("pdf", pdf_bytes)
        total_pages = doc.page_count
        doc.close()

        thumbnails = []
        start_render = time.time()

        for i in range(total_pages):
            thumbnail = render_pdf_page_to_image(pdf_bytes, i, target_width=480)
            thumbnails.append(thumbnail)
            render_progress = 50 + int((i + 1) / total_pages * 45)
            yield f"data: {json.dumps({'stage': 'rendering', 'progress': render_progress, 'message': f'rendered {i+1}/{total_pages}', 'current_slide': i+1, 'total_slides': total_pages})}\n\n"
            await asyncio.sleep(0.01)

        render_time = time.time() - start_render
        total_time = convert_time + render_time

        yield f"data: {json.dumps({'stage': 'done', 'progress': 100, 'message': f'completed in {total_time:.1f}s', 'thumbnails': thumbnails, 'format': 'jpeg'})}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive", "X-Accel-Buffering": "no"},
    )


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8004)
