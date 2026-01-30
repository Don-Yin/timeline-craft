"""thumbnail generation endpoints"""
import io
import logging
import time

from fastapi import APIRouter, HTTPException
from pptx import Presentation

from .schemas import ThumbnailResponse, SlideCountResponse, AllThumbnailsResponse
from ..modules.cache import PdfCache
from ..modules.storage import get_file_from_minio, get_cached_pdf, store_cached_pdf, MINIO_BUCKET_UPLOADS
from ..modules.pdf_renderer import convert_pptx_to_pdf_bytes, render_slide_thumbnail, render_all_thumbnails

logger = logging.getLogger(__name__)
router = APIRouter(tags=["thumbnails"])

pdf_cache = PdfCache(ttl_seconds=300)


def get_or_create_pdf(file_id: str) -> bytes:
    """get pdf from cache or convert pptx to pdf"""
    pdf_bytes = get_cached_pdf(file_id)
    if pdf_bytes:
        return pdf_bytes

    pdf_bytes = pdf_cache.get(file_id)
    if pdf_bytes:
        return pdf_bytes

    logger.info(f"no cached pdf for {file_id}, converting via libreoffice (~28s)")
    file_buffer = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    pptx_bytes = file_buffer.getvalue()
    pdf_bytes = convert_pptx_to_pdf_bytes(file_id, pptx_bytes)

    pdf_cache.set(file_id, pdf_bytes)
    store_cached_pdf(file_id, pdf_bytes)
    return pdf_bytes


@router.get("/get-thumbnail/{file_id}/{slide_index}", response_model=ThumbnailResponse, summary="Get slide thumbnail")
async def get_thumbnail(file_id: str, slide_index: int):
    """generate a png thumbnail for a specific slide in a presentation"""
    logger.info(f"fetching thumbnail for file {file_id}, slide {slide_index}")

    file_buffer = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    pptx_bytes = file_buffer.getvalue()
    prs = Presentation(io.BytesIO(pptx_bytes))

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise HTTPException(status_code=404, detail="slide index out of range")

    pdf_bytes = get_or_create_pdf(file_id)
    img_str = render_slide_thumbnail(pdf_bytes, slide_index)
    return ThumbnailResponse(slide_index=slide_index, image_base64=img_str)


@router.get("/get-all-thumbnails/{file_id}", response_model=AllThumbnailsResponse, summary="Get all thumbnails at once")
async def get_all_thumbnails(file_id: str):
    """generate all slide thumbnails - uses persistent pdf cache for speed (~0.8s if cached)"""
    logger.info(f"fetching all thumbnails for file {file_id}")
    start_time = time.time()

    pdf_bytes = get_or_create_pdf(file_id)
    thumbnails = render_all_thumbnails(pdf_bytes)

    elapsed = time.time() - start_time
    logger.info(f"generated {len(thumbnails)} thumbnails for {file_id} in {elapsed:.2f}s")
    return AllThumbnailsResponse(file_id=file_id, thumbnails=thumbnails)


@router.get("/get-slide-count/{file_id}", response_model=SlideCountResponse, summary="Get slide count")
async def get_slide_count(file_id: str):
    """get the total number of slides in a presentation"""
    file_data = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    prs = Presentation(file_data)
    return SlideCountResponse(file_id=file_id, slide_count=len(prs.slides))


@router.post("/pre-generate-pdf/{file_id}", summary="Pre-generate PDF cache")
async def pre_generate_pdf(file_id: str):
    """pre-convert pptx to pdf and cache in minio for instant thumbnail loading later"""
    logger.info(f"pre-generating pdf cache for {file_id}")

    existing = get_cached_pdf(file_id)
    if existing:
        logger.info(f"pdf already cached for {file_id}")
        return {"status": "already_cached", "file_id": file_id}

    start_time = time.time()
    file_buffer = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    pptx_bytes = file_buffer.getvalue()

    pdf_bytes = convert_pptx_to_pdf_bytes(file_id, pptx_bytes)
    store_cached_pdf(file_id, pdf_bytes)
    pdf_cache.set(file_id, pdf_bytes)

    elapsed = time.time() - start_time
    logger.info(f"pre-generated pdf for {file_id} in {elapsed:.1f}s")
    return {"status": "generated", "file_id": file_id, "time_seconds": round(elapsed, 1)}

