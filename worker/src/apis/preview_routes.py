"""preview generation endpoints with sidebar applied"""

import io
import logging

from fastapi import APIRouter, HTTPException
from pptx import Presentation
from pptx.dml.color import RGBColor
from webcolors import hex_to_rgb

from .schemas import ThumbnailResponse, PreviewRequest, AllPreviewsResponse
from ..modules.cache import PdfCache
from ..modules.storage import get_file_from_minio, MINIO_BUCKET_UPLOADS
from ..modules.pdf_renderer import convert_pptx_to_pdf_bytes, render_slide_thumbnail, render_all_thumbnails
from timeline import set_sidebar_timeline, move_elements_to_right, Configurations

logger = logging.getLogger(__name__)
router = APIRouter(tags=["thumbnails"])

preview_cache = PdfCache(ttl_seconds=60, max_entries=10)


def build_preview_config(request: PreviewRequest) -> Configurations:
    """build configurations object from preview request"""
    return Configurations(
        sidebar_width=request.sidebar_width,
        sidebar_item_height=request.sidebar_item_height,
        sidebar_color=RGBColor(*hex_to_rgb(request.sidebar_color_hex)),
        indicator_color=RGBColor(*hex_to_rgb(request.indicator_color_hex)),
        sidebar_item_font_color=RGBColor(*hex_to_rgb(request.sidebar_item_font_color_hex)),
        sidebar_transparency=request.sidebar_transparency * 1000,
    )


def build_cache_key(file_id: str, request: PreviewRequest) -> str:
    """build unique cache key for preview"""
    return f"{file_id}:{request.sidebar_width}:{request.sidebar_item_height}:{request.sidebar_color_hex}:{request.indicator_color_hex}:{request.sidebar_item_font_color_hex}:{request.sidebar_transparency}:{hash(tuple(request.tags))}"


def generate_preview_pdf(file_id: str, prs: Presentation, request: PreviewRequest) -> bytes:
    """apply sidebar and convert to pdf"""
    config = build_preview_config(request)
    move_elements_to_right(prs, config=config)
    set_sidebar_timeline(ppt=prs, tags=request.tags, config=config)

    preview_buffer = io.BytesIO()
    prs.save(preview_buffer)
    preview_buffer.seek(0)

    return convert_pptx_to_pdf_bytes(f"preview-{file_id}", preview_buffer.getvalue())


@router.post(
    "/get-preview-thumbnail/{file_id}/{slide_index}", response_model=ThumbnailResponse, summary="Get preview with sidebar"
)
async def get_preview_thumbnail(file_id: str, slide_index: int, request: PreviewRequest):
    """generate a preview thumbnail showing how the slide will look with sidebar applied"""
    logger.info(f"generating preview for file {file_id}, slide {slide_index}")

    file_buffer = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    pptx_bytes = file_buffer.getvalue()
    prs = Presentation(io.BytesIO(pptx_bytes))

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise HTTPException(status_code=404, detail="slide index out of range")

    if len(request.tags) != len(prs.slides):
        raise HTTPException(status_code=400, detail=f"tags count ({len(request.tags)}) must match slides ({len(prs.slides)})")

    cache_key = build_cache_key(file_id, request)
    pdf_bytes = preview_cache.get(cache_key)

    if pdf_bytes is None:
        logger.info(f"generating preview pdf for {file_id}")
        pdf_bytes = generate_preview_pdf(file_id, prs, request)
        preview_cache.set(cache_key, pdf_bytes)

    img_str = render_slide_thumbnail(pdf_bytes, slide_index)
    return ThumbnailResponse(slide_index=slide_index, image_base64=img_str)


@router.post(
    "/get-all-preview-thumbnails/{file_id}", response_model=AllPreviewsResponse, summary="Get all previews with sidebar"
)
async def get_all_preview_thumbnails(file_id: str, request: PreviewRequest):
    """generate all preview thumbnails with sidebar applied in a single request"""
    logger.info(f"generating all preview thumbnails for file {file_id}")

    file_buffer = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    pptx_bytes = file_buffer.getvalue()
    prs = Presentation(io.BytesIO(pptx_bytes))

    if len(request.tags) != len(prs.slides):
        raise HTTPException(status_code=400, detail=f"tags count ({len(request.tags)}) must match slides ({len(prs.slides)})")

    cache_key = build_cache_key(file_id, request)
    pdf_bytes = preview_cache.get(cache_key)

    if pdf_bytes is None:
        logger.info(f"generating preview pdf for {file_id}")
        pdf_bytes = generate_preview_pdf(file_id, prs, request)
        preview_cache.set(cache_key, pdf_bytes)

    thumbnails = render_all_thumbnails(pdf_bytes, target_width=320, use_jpeg=True)
    logger.info(f"generated {len(thumbnails)} preview thumbnails for {file_id}")
    return AllPreviewsResponse(file_id=file_id, thumbnails=thumbnails, format="jpeg")


def get_preview_cache() -> PdfCache:
    """expose preview cache for cleanup from other modules"""
    return preview_cache
