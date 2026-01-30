import logging, io, json, uuid, time
from threading import Lock

from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation

from .schemas import ThumbnailResponse, ProcessRequest, ProcessResponse, ListProcessedResponse, SlideCountResponse, PreviewRequest, AllThumbnailsResponse, AllPreviewsResponse
from ..modules.cache import PdfCache
from ..modules.storage import (
    get_file_from_minio,
    upload_file_to_minio,
    list_files_in_bucket,
    MINIO_BUCKET_UPLOADS,
    MINIO_BUCKET_PROCESSED,
)
from ..modules.pdf_renderer import convert_pptx_to_pdf_bytes, render_slide_thumbnail, render_all_thumbnails
from ..modules.pptx_processor import process_presentation
from ..modules.progressive_processor import ProgressiveProcessor
from timeline import set_sidebar_timeline, move_elements_to_right, Configurations

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

router = APIRouter()

pdf_cache = PdfCache(ttl_seconds=300)
preview_cache = PdfCache(ttl_seconds=60, max_entries=10)

processed_files_store: dict[str, tuple[float, io.BytesIO]] = {}
store_lock = Lock()
STORE_TTL = 300


@router.get("/get-thumbnail/{file_id}/{slide_index}", response_model=ThumbnailResponse, tags=["thumbnails"], summary="Get slide thumbnail")
async def get_thumbnail(file_id: str, slide_index: int):
    """generate a png thumbnail for a specific slide in a presentation"""
    logger.info(f"fetching thumbnail for file {file_id}, slide {slide_index}")

    file_buffer = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    pptx_bytes = file_buffer.getvalue()
    prs = Presentation(io.BytesIO(pptx_bytes))

    logger.info(f"loaded presentation with {len(prs.slides)} slides")

    if slide_index < 0 or slide_index >= len(prs.slides):
        logger.warning(f"slide index {slide_index} out of range")
        raise HTTPException(status_code=404, detail="slide index out of range")

    pdf_bytes = pdf_cache.get(file_id)
    if pdf_bytes is None:
        logger.info(f"no cached pdf for {file_id}, converting via libreoffice")
        pdf_bytes = convert_pptx_to_pdf_bytes(file_id, pptx_bytes)
        pdf_cache.set(file_id, pdf_bytes)

    img_str = render_slide_thumbnail(pdf_bytes, slide_index)
    return ThumbnailResponse(slide_index=slide_index, image_base64=img_str)


@router.get("/get-all-thumbnails/{file_id}", response_model=AllThumbnailsResponse, tags=["thumbnails"], summary="Get all thumbnails at once")
async def get_all_thumbnails(file_id: str):
    """generate all slide thumbnails in a single request for faster loading"""
    logger.info(f"fetching all thumbnails for file {file_id}")

    file_buffer = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    pptx_bytes = file_buffer.getvalue()

    pdf_bytes = pdf_cache.get(file_id)
    if pdf_bytes is None:
        logger.info(f"no cached pdf for {file_id}, converting via libreoffice")
        pdf_bytes = convert_pptx_to_pdf_bytes(file_id, pptx_bytes)
        pdf_cache.set(file_id, pdf_bytes)

    thumbnails = render_all_thumbnails(pdf_bytes)
    logger.info(f"generated {len(thumbnails)} thumbnails for {file_id}")
    return AllThumbnailsResponse(file_id=file_id, thumbnails=thumbnails)


@router.get("/get-slide-count/{file_id}", response_model=SlideCountResponse, tags=["thumbnails"], summary="Get slide count")
async def get_slide_count(file_id: str):
    """get the total number of slides in a presentation"""
    file_data = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    prs = Presentation(file_data)
    return SlideCountResponse(file_id=file_id, slide_count=len(prs.slides))


@router.post("/get-preview-thumbnail/{file_id}/{slide_index}", response_model=ThumbnailResponse, tags=["thumbnails"], summary="Get preview with sidebar")
async def get_preview_thumbnail(file_id: str, slide_index: int, request: PreviewRequest):
    """generate a preview thumbnail showing how the slide will look with sidebar applied"""
    logger.info(f"generating preview for file {file_id}, slide {slide_index}")

    file_buffer = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    pptx_bytes = file_buffer.getvalue()
    prs = Presentation(io.BytesIO(pptx_bytes))

    if slide_index < 0 or slide_index >= len(prs.slides):
        raise HTTPException(status_code=404, detail="slide index out of range")

    if len(request.tags) != len(prs.slides):
        raise HTTPException(status_code=400, detail=f"tags count ({len(request.tags)}) must match slides ({len(prs.slides)})")

    cache_key = f"{file_id}:{request.sidebar_width}:{request.sidebar_item_height}:{hash(tuple(request.tags))}"
    pdf_bytes = preview_cache.get(cache_key)

    if pdf_bytes is None:
        logger.info(f"generating preview pdf for {file_id}")
        config = Configurations(sidebar_width=request.sidebar_width, sidebar_item_height=request.sidebar_item_height)
        move_elements_to_right(prs, config=config)
        set_sidebar_timeline(ppt=prs, tags=request.tags, config=config)

        preview_buffer = io.BytesIO()
        prs.save(preview_buffer)
        preview_buffer.seek(0)

        pdf_bytes = convert_pptx_to_pdf_bytes(f"preview-{file_id}", preview_buffer.getvalue())
        preview_cache.set(cache_key, pdf_bytes)

    img_str = render_slide_thumbnail(pdf_bytes, slide_index)
    return ThumbnailResponse(slide_index=slide_index, image_base64=img_str)


@router.post("/get-all-preview-thumbnails/{file_id}", response_model=AllPreviewsResponse, tags=["thumbnails"], summary="Get all previews with sidebar")
async def get_all_preview_thumbnails(file_id: str, request: PreviewRequest):
    """generate all preview thumbnails with sidebar applied in a single request"""
    logger.info(f"generating all preview thumbnails for file {file_id}")

    file_buffer = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    pptx_bytes = file_buffer.getvalue()
    prs = Presentation(io.BytesIO(pptx_bytes))

    if len(request.tags) != len(prs.slides):
        raise HTTPException(status_code=400, detail=f"tags count ({len(request.tags)}) must match slides ({len(prs.slides)})")

    cache_key = f"{file_id}:{request.sidebar_width}:{request.sidebar_item_height}:{hash(tuple(request.tags))}"
    pdf_bytes = preview_cache.get(cache_key)

    if pdf_bytes is None:
        logger.info(f"generating preview pdf for {file_id}")
        config = Configurations(sidebar_width=request.sidebar_width, sidebar_item_height=request.sidebar_item_height)
        move_elements_to_right(prs, config=config)
        set_sidebar_timeline(ppt=prs, tags=request.tags, config=config)

        preview_buffer = io.BytesIO()
        prs.save(preview_buffer)
        preview_buffer.seek(0)

        pdf_bytes = convert_pptx_to_pdf_bytes(f"preview-{file_id}", preview_buffer.getvalue())
        preview_cache.set(cache_key, pdf_bytes)

    # Use lower resolution (320px) and JPEG for faster preview generation
    thumbnails = render_all_thumbnails(pdf_bytes, target_width=320, use_jpeg=True)
    logger.info(f"generated {len(thumbnails)} preview thumbnails for {file_id}")
    return AllPreviewsResponse(file_id=file_id, thumbnails=thumbnails, format="jpeg")


@router.post("/process-file/{file_id}", response_model=ProcessResponse, tags=["processing"], summary="Process and store presentation")
async def process_file(file_id: str, request: ProcessRequest):
    """process a presentation with timeline sidebar and store the result"""
    file_data = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    processed_buffer = process_presentation(file_data, request)
    upload_file_to_minio(
        MINIO_BUCKET_PROCESSED,
        file_id,
        processed_buffer,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    return ProcessResponse(file_id=file_id, processed_file_id=file_id, message="file processed successfully")


@router.get("/list-processed", response_model=ListProcessedResponse, tags=["processing"], summary="List processed files")
async def list_processed():
    """get a list of all processed file ids"""
    files = list_files_in_bucket(MINIO_BUCKET_PROCESSED)
    return ListProcessedResponse(files=files)


@router.post("/process-and-download/{file_id}", tags=["processing"], summary="Process and download directly")
async def process_and_download(file_id: str, request: ProcessRequest):
    """process a presentation and return it directly as a downloadable file"""
    logger.info(f"processing file {file_id} for download")
    file_data = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    processed_buffer = process_presentation(file_data, request)

    return StreamingResponse(
        processed_buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename=timeline-{file_id}.pptx"},
    )


def _cleanup_expired_files():
    """remove expired files from temporary store"""
    now = time.time()
    with store_lock:
        expired = [k for k, (ts, _) in processed_files_store.items() if now - ts > STORE_TTL]
        for k in expired:
            del processed_files_store[k]


def _generate_progress_stream(file_id: str, request: ProcessRequest):
    """generator for sse progress events"""
    _cleanup_expired_files()

    file_data = get_file_from_minio(MINIO_BUCKET_UPLOADS, file_id)
    processor = ProgressiveProcessor(file_data, request)

    result_buffer = None

    for progress_event in processor.process_with_progress():
        if isinstance(progress_event, io.BytesIO):
            result_buffer = progress_event
            break

        event_data = json.dumps(progress_event)
        yield f"data: {event_data}\n\n"

        if progress_event.get("stage") == "complete":
            break

    if result_buffer is None:
        result_buffer = io.BytesIO()
        processor.prs.save(result_buffer)
        result_buffer.seek(0)

    job_id = str(uuid.uuid4())
    with store_lock:
        processed_files_store[job_id] = (time.time(), result_buffer)

    done_event = json.dumps({"stage": "ready", "progress": 100, "job_id": job_id})
    yield f"data: {done_event}\n\n"


@router.post("/process-with-progress/{file_id}", tags=["processing"], summary="Process with real-time progress")
async def process_with_progress(file_id: str, request: ProcessRequest):
    """process a presentation with sse for real-time progress updates"""
    logger.info(f"starting progressive processing for {file_id}")

    return StreamingResponse(
        _generate_progress_stream(file_id, request), media_type="text/event-stream", headers={"Cache-Control": "no-cache", "Connection": "keep-alive", "X-Accel-Buffering": "no"}
    )


@router.get("/download-processed/{job_id}", tags=["processing"], summary="Download processed file by job ID")
async def download_processed(job_id: str):
    """download a processed file using the job id from process-with-progress"""
    _cleanup_expired_files()

    with store_lock:
        entry = processed_files_store.get(job_id)
        if not entry:
            raise HTTPException(status_code=404, detail="processed file not found or expired")

        _, buffer = entry
        buffer.seek(0)
        content = buffer.read()
        del processed_files_store[job_id]

    return StreamingResponse(
        io.BytesIO(content),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename=timeline-{job_id}.pptx"},
    )
